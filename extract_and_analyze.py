import zipfile
import os
from pptx import Presentation
import json

def extract_zip(zip_path, extract_to='./neonet_project'):
    """Extract ZIP archive"""
    print(f"Extracting {zip_path}...")
    with zipfile.ZipFile(zip_path, 'r') as zip_ref:
        zip_ref.extractall(extract_to)
    print(f"Extracted to {extract_to}")
    return extract_to

def parse_pptx(pptx_path):
    """Parse PowerPoint presentation and extract text content"""
    print(f"\n{'='*80}")
    print(f"Analyzing: {os.path.basename(pptx_path)}")
    print(f"{'='*80}\n")
    
    prs = Presentation(pptx_path)
    slides_data = []
    
    for slide_num, slide in enumerate(prs.slides, 1):
        slide_content = {
            'slide_number': slide_num,
            'title': '',
            'content': []
        }
        
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                text = shape.text.strip()
                if shape.shape_type == 14:
                    slide_content['title'] = text
                    print(f"Slide {slide_num}: {text}")
                    print("-" * 80)
                else:
                    slide_content['content'].append(text)
                    print(text)
        
        print("\n")
        slides_data.append(slide_content)
    
    return slides_data

def analyze_project_structure(project_dir):
    """Analyze extracted project structure"""
    print(f"\n{'='*80}")
    print("PROJECT STRUCTURE ANALYSIS")
    print(f"{'='*80}\n")
    
    for root, dirs, files in os.walk(project_dir):
        level = root.replace(project_dir, '').count(os.sep)
        indent = ' ' * 2 * level
        print(f'{indent}{os.path.basename(root)}/')
        subindent = ' ' * 2 * (level + 1)
        for file in files:
            print(f'{subindent}{file}')

def main():
    zip_file = 'attached_assets/NeoNetstage8pqcinteropzip_1763535779243.zip'
    pptx_files = [
        'attached_assets/NeoNet_PitchDeck_Extended_1763535787630.pptx',
        'attached_assets/NeoNet_PitchDeck_Compact_English_1763535793221.pptx'
    ]
    
    extract_dir = extract_zip(zip_file)
    
    all_presentation_data = {}
    for pptx_file in pptx_files:
        if os.path.exists(pptx_file):
            presentation_name = os.path.basename(pptx_file)
            all_presentation_data[presentation_name] = parse_pptx(pptx_file)
    
    analyze_project_structure(extract_dir)
    
    with open('presentation_analysis.json', 'w', encoding='utf-8') as f:
        json.dump(all_presentation_data, f, indent=2, ensure_ascii=False)
    print("\nPresentation data saved to presentation_analysis.json")

if __name__ == '__main__':
    main()
